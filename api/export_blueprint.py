from flask import Blueprint, request, jsonify, send_file
import base64
import io
import os
import tempfile
import zipfile
from datetime import datetime
import uuid

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Image, PageBreak
from reportlab.lib.utils import ImageReader

# Image processing
from PIL import Image as PILImage

# Create blueprint
export_bp = Blueprint('export', __name__, url_prefix='/api/export')

# Store temporary sessions (in production, use Redis or database)
export_sessions = {}


class ExportProcessor:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()

    def decode_base64_image(self, base64_string):
        """Convert base64 string to PIL Image"""
        # Remove data URL prefix if present
        if base64_string.startswith('data:image'):
            base64_string = base64_string.split(',')[1]

        image_data = base64.b64decode(base64_string)
        return PILImage.open(io.BytesIO(image_data))

    def create_powerpoint(self, images, metadata=None):
        """Create PowerPoint presentation from images"""
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for i, img_data in enumerate(images):
            # Create slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

            # Convert base64 to image
            pil_image = self.decode_base64_image(img_data)

            # Save image to temporary file
            temp_img_path = os.path.join(self.temp_dir, f'slide_{i}_{uuid.uuid4().hex[:8]}.png')
            pil_image.save(temp_img_path, 'PNG')

            # Add image to slide (fill entire slide)
            slide.shapes.add_picture(
                temp_img_path,
                Inches(0),
                Inches(0),
                width=Inches(13.33),
                height=Inches(7.5)
            )

            # Clean up temp image
            try:
                os.remove(temp_img_path)
            except:
                pass  # Ignore cleanup errors

        # Save PowerPoint to memory
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def create_pdf(self, images, metadata=None):
        """Create PDF from images"""
        pdf_buffer = io.BytesIO()

        # Create PDF document (landscape orientation for better fit)
        doc = SimpleDocTemplate(
            pdf_buffer,
            pagesize=(A4[1], A4[0]),  # Landscape
            rightMargin=0,
            leftMargin=0,
            topMargin=0,
            bottomMargin=0
        )

        story = []

        for i, img_data in enumerate(images):
            # Convert base64 to image
            pil_image = self.decode_base64_image(img_data)

            # Convert to RGB if necessary
            if pil_image.mode != 'RGB':
                pil_image = pil_image.convert('RGB')

            # Save to temporary buffer
            img_buffer = io.BytesIO()
            pil_image.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # Add image to PDF (landscape)
            img = Image(img_buffer, width=A4[1], height=A4[0])
            story.append(img)

            # Add page break except for last image
            if i < len(images) - 1:
                story.append(PageBreak())

        # Build PDF
        doc.build(story)
        pdf_buffer.seek(0)

        return pdf_buffer

    def create_zip(self, images, metadata=None):
        """Create ZIP file with images"""
        zip_buffer = io.BytesIO()

        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            # Add README
            total_slides = len(images)
            created_at = metadata.get('created_at',
                                      datetime.now().isoformat()) if metadata else datetime.now().isoformat()

            readme_content = f"""Animated Circle Presentation
Generated: {created_at}
Total slides: {total_slides}

Instructions:
- Import these images into your presentation software
- Each image represents one animation frame
- Recommended slide timing: 2 seconds per slide

Files included:
"""

            # Add images
            for i, img_data in enumerate(images):
                pil_image = self.decode_base64_image(img_data)

                # Save image to buffer
                img_buffer = io.BytesIO()
                pil_image.save(img_buffer, format='PNG', optimize=True)
                img_buffer.seek(0)

                # Add to ZIP
                filename = f"slide_{i + 1:02d}.png"
                zip_file.writestr(filename, img_buffer.getvalue())
                readme_content += f"- {filename}\n"

            zip_file.writestr("README.txt", readme_content)

        zip_buffer.seek(0)
        return zip_buffer


# Initialize processor
processor = ExportProcessor()


@export_bp.route('/', methods=['POST'])
def export_presentation():
    """Main export endpoint - handles single request with all images"""
    try:
        data = request.get_json()
        images = data.get('images', [])
        export_type = data.get('export_type', 'pdf')
        metadata = data.get('metadata', {})

        if not images:
            return jsonify({'error': 'No images provided'}), 400

        if len(images) > 50:  # Reasonable limit
            return jsonify({'error': 'Too many images (max 50)'}), 400

        # Generate file based on type
        try:
            if export_type == 'pptx':
                file_buffer = processor.create_powerpoint(images, metadata)
                mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                filename = 'animated_presentation.pptx'
            elif export_type == 'pdf':
                file_buffer = processor.create_pdf(images, metadata)
                mimetype = 'application/pdf'
                filename = 'animated_presentation.pdf'
            elif export_type == 'zip':
                file_buffer = processor.create_zip(images, metadata)
                mimetype = 'application/zip'
                filename = 'animated_slides.zip'
            else:
                return jsonify({'error': 'Invalid export type. Use: pptx, pdf, or zip'}), 400

            return send_file(
                file_buffer,
                mimetype=mimetype,
                as_attachment=True,
                download_name=filename
            )

        except Exception as processing_error:
            print(f"Processing error: {str(processing_error)}")
            return jsonify({'error': f'Processing failed: {str(processing_error)}'}), 500

    except Exception as e:
        print(f"Export error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@export_bp.route('/session', methods=['POST'])
def create_export_session():
    """Create a new export session for progressive upload"""
    try:
        data = request.get_json()
        session_id = str(uuid.uuid4())

        export_sessions[session_id] = {
            'export_type': data.get('export_type', 'pdf'),
            'total_images': data.get('total_images', 0),
            'images': {},
            'metadata': data.get('metadata', {}),
            'created_at': datetime.now(),
            'expires_at': datetime.now().timestamp() + 3600  # 1 hour expiry
        }

        return jsonify({
            'session_id': session_id,
            'expires_in': 3600,
            'max_images': 50
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@export_bp.route('/upload', methods=['POST'])
def upload_image():
    """Upload single image to session"""
    try:
        data = request.get_json()
        session_id = data.get('session_id')
        image_index = data.get('image_index')
        image_data = data.get('image_data')

        if not all([session_id, image_index is not None, image_data]):
            return jsonify({'error': 'Missing required fields'}), 400

        if session_id not in export_sessions:
            return jsonify({'error': 'Invalid or expired session'}), 400

        session_data = export_sessions[session_id]

        # Check if session expired
        if datetime.now().timestamp() > session_data['expires_at']:
            del export_sessions[session_id]
            return jsonify({'error': 'Session expired'}), 400

        # Validate image index
        if image_index >= session_data['total_images'] or image_index < 0:
            return jsonify({'error': 'Invalid image index'}), 400

        # Store image
        session_data['images'][image_index] = image_data

        return jsonify({
            'status': 'uploaded',
            'images_received': len(session_data['images']),
            'total_expected': session_data['total_images']
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@export_bp.route('/generate', methods=['POST'])
def generate_from_session():
    """Generate file from session data"""
    try:
        data = request.get_json()
        session_id = data.get('session_id')

        if not session_id:
            return jsonify({'error': 'Missing session_id'}), 400

        if session_id not in export_sessions:
            return jsonify({'error': 'Invalid or expired session'}), 400

        session_data = export_sessions[session_id]

        # Check if we have all images
        expected_images = session_data['total_images']
        received_images = len(session_data['images'])

        if received_images != expected_images:
            return jsonify({
                'error': f'Missing images. Expected {expected_images}, got {received_images}'
            }), 400

        # Convert images dict to ordered list
        images = []
        for i in range(expected_images):
            if i not in session_data['images']:
                return jsonify({'error': f'Missing image at index {i}'}), 400
            images.append(session_data['images'][i])

        export_type = session_data['export_type']
        metadata = session_data['metadata']

        # Generate file
        try:
            if export_type == 'pptx':
                file_buffer = processor.create_powerpoint(images, metadata)
                mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                filename = 'animated_presentation.pptx'
            elif export_type == 'pdf':
                file_buffer = processor.create_pdf(images, metadata)
                mimetype = 'application/pdf'
                filename = 'animated_presentation.pdf'
            elif export_type == 'zip':
                file_buffer = processor.create_zip(images, metadata)
                mimetype = 'application/zip'
                filename = 'animated_slides.zip'
            else:
                return jsonify({'error': 'Invalid export type'}), 400

            # Clean up session after successful generation
            del export_sessions[session_id]

            return send_file(
                file_buffer,
                mimetype=mimetype,
                as_attachment=True,
                download_name=filename
            )

        except Exception as processing_error:
            return jsonify({'error': f'Generation failed: {str(processing_error)}'}), 500

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@export_bp.route('/health', methods=['GET'])
def export_health():
    """Health check for export service"""
    active_sessions = len(export_sessions)

    # Clean up expired sessions
    current_time = datetime.now().timestamp()
    expired_sessions = [sid for sid, data in export_sessions.items()
                        if current_time > data['expires_at']]

    for sid in expired_sessions:
        del export_sessions[sid]

    return jsonify({
        'status': 'healthy',
        'active_sessions': active_sessions,
        'cleaned_expired': len(expired_sessions),
        'timestamp': datetime.now().isoformat()
    })


@export_bp.route('/sessions', methods=['GET'])
def list_sessions():
    """Debug endpoint to list active sessions"""
    sessions_info = []
    current_time = datetime.now().timestamp()

    for sid, data in export_sessions.items():
        sessions_info.append({
            'session_id': sid,
            'export_type': data['export_type'],
            'total_images': data['total_images'],
            'images_received': len(data['images']),
            'created_at': data['created_at'].isoformat(),
            'expires_in': max(0, int(data['expires_at'] - current_time))
        })

    return jsonify({
        'sessions': sessions_info,
        'count': len(sessions_info)
    })


# Cleanup function (call this periodically)
def cleanup_expired_sessions():
    """Remove expired sessions"""
    current_time = datetime.now().timestamp()
    expired_sessions = [sid for sid, data in export_sessions.items()
                        if current_time > data['expires_at']]

    for sid in expired_sessions:
        del export_sessions[sid]

    return len(expired_sessions)